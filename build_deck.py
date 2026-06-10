# -*- coding: utf-8 -*-
"""Build DOS Curacao presentation deck (.pptx) from the website content."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(ROOT, "assets", "images")
LOGO = os.path.join(ROOT, "assets", "logos")

# ---- brand ----
BLUE   = RGBColor(0x08, 0x52, 0xA7)
YELLOW = RGBColor(0xF8, 0xCB, 0x0C)
RED    = RGBColor(0xD7, 0x20, 0x24)
INK    = RGBColor(0x0A, 0x1A, 0x2F)
SLATE  = RGBColor(0x5A, 0x6B, 0x82)
WASH   = RGBColor(0xE4, 0xED, 0xFB)
CLOUD  = RGBColor(0xEA, 0xF0, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Archivo"
MONO = "JetBrains Mono"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- image dimensions (no PIL dependency) ----------
def img_size(path):
    with open(path, "rb") as f:
        d = f.read()
    if d[:8] == b"\x89PNG\r\n\x1a\n":
        return int.from_bytes(d[16:20], "big"), int.from_bytes(d[20:24], "big")
    if d[:2] == b"\xff\xd8":
        i = 2
        while i < len(d):
            if d[i] != 0xFF:
                i += 1; continue
            m = d[i+1]
            if m in (0xC0, 0xC1, 0xC2, 0xC3):
                return int.from_bytes(d[i+7:i+9], "big"), int.from_bytes(d[i+5:i+7], "big")
            i += 2 + int.from_bytes(d[i+2:i+4], "big")
    return (16, 9)

# ---------- helpers ----------
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def rect(s, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=None):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def set_alpha(shape, opacity_pct):
    """opacity_pct: 0..100 (100 = opaque)."""
    sf = shape._element.spPr.find(qn('a:solidFill'))
    srgb = sf.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity_pct * 1000))})
    srgb.append(a)

def txt(s, l, t, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        font=FONT, anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = runs if isinstance(runs, list) else [runs]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = line
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
    return tb

def bullets(s, l, t, w, h, items, size=15, color=INK, bcolor=BLUE, gap=8, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        rb = p.add_run(); rb.text = "■  "
        rb.font.size = Pt(size-4); rb.font.color.rgb = bcolor; rb.font.bold = True; rb.font.name = FONT
        if isinstance(it, tuple):
            r1 = p.add_run(); r1.text = it[0] + "  "
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = FONT
            r2 = p.add_run(); r2.text = it[1]
            r2.font.size = Pt(size); r2.font.color.rgb = SLATE; r2.font.name = FONT
        else:
            r1 = p.add_run(); r1.text = it
            r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.name = FONT
    return tb

def cover_img(s, path, l, t, w, h):
    """Insert picture filling box l,t,w,h with center-crop (no distortion)."""
    iw, ih = img_size(path)
    ar = iw / ih; box_ar = w / h
    pic = s.shapes.add_picture(path, l, t, width=w, height=h)
    if ar > box_ar:
        crop = (1 - box_ar / ar) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - ar / box_ar) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic

def contain_img(s, path, l, t, w, h):
    iw, ih = img_size(path); ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = Emu(int(w / ar))
    else:
        nh = h; nw = Emu(int(h * ar))
    nl = l + (w - nw) // 2; nt = t + (h - nh) // 2
    return s.shapes.add_picture(path, nl, nt, width=nw, height=nh)

def tri_bar(s, l, t, w, h):
    seg = w // 3
    rect(s, l, t, seg, h, BLUE)
    rect(s, l + seg, t, seg, h, YELLOW)
    rect(s, l + 2*seg, t, w - 2*seg, h, RED)

PAGE = [0]
def header(s, num, eyebrow, ecolor, title, title_size=34, dark=False, title_w=None):
    if title_w is None:
        title_w = Inches(11.9)
    tcol = WHITE if dark else INK
    # number
    txt(s, Inches(0.7), Inches(0.55), Inches(2), Inches(0.5),
        f"{num:02d}", size=15, color=ecolor, bold=True, font=MONO)
    # eyebrow with dot bar
    rect(s, Inches(0.7), Inches(1.02), Inches(0.32), Inches(0.045), ecolor)
    txt(s, Inches(1.12), Inches(0.83), Inches(9), Inches(0.4),
        eyebrow.upper(), size=12, color=ecolor, bold=True, font=MONO)
    # title
    txt(s, Inches(0.7), Inches(1.35), title_w, Inches(1.9),
        title, size=title_size, color=tcol, bold=True, font=FONT, spacing=0.98)

def footer(s, dark=False):
    PAGE[0] += 1
    col = RGBColor(0x9F,0xB3,0xCC) if dark else SLATE
    tri_bar(s, Inches(0.7), Inches(7.16), Inches(0.55), Inches(0.07))
    txt(s, Inches(1.4), Inches(7.0), Inches(7), Inches(0.4),
        "DOS Curacao  ·  Workforce Management Platform", size=9, color=col, font=MONO)
    txt(s, Inches(11.0), Inches(7.0), Inches(1.63), Inches(0.4),
        f"{PAGE[0]:02d}", size=9, color=col, font=MONO, align=PP_ALIGN.RIGHT)

def card(s, l, t, w, h, fill, line=CLOUD):
    return rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=line, line_w=Pt(1))

# =========================================================
# SLIDE 1 — TITLE
# =========================================================
s = slide(); bg(s, WHITE)
rect(s, 0, 0, SW, Inches(0.18), BLUE)
tri_bar(s, 0, 0, SW, Inches(0.18))
crest = os.path.join(LOGO, "dos-crest-hd.png")
if os.path.exists(crest):
    contain_img(s, crest, Inches(5.67), Inches(1.15), Inches(2.0), Inches(1.5))
txt(s, Inches(1), Inches(2.85), Inches(11.33), Inches(0.4),
    "DIENST OPENBARE SCHOLEN  ·  CURAÇAO", size=13, color=BLUE, bold=True,
    font=MONO, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.35), Inches(11.33), Inches(1.6),
    "The Future of School\nWorkforce Management.", size=46, color=INK, bold=True,
    font=FONT, align=PP_ALIGN.CENTER, spacing=0.98)
txt(s, Inches(2.4), Inches(5.25), Inches(8.53), Inches(1.1),
    "A secure, mobile-first attendance platform that gives every teacher a simple way to check in — and gives school leadership real-time visibility across Curaçao's public schools.",
    size=15, color=SLATE, font=FONT, align=PP_ALIGN.CENTER, spacing=1.1)
txt(s, Inches(1), Inches(6.7), Inches(11.33), Inches(0.4),
    "An official platform of the Ministerie van Onderwijs, Wetenschap, Cultuur & Sport",
    size=10.5, color=SLATE, font=MONO, align=PP_ALIGN.CENTER)

# =========================================================
# Reusable content layout: left text, right cover image
# =========================================================
TXT_W = Inches(6.5)
IMG_L = Inches(7.55); IMG_T = Inches(1.35); IMG_W = Inches(5.08); IMG_H = Inches(5.3)

def img_panel(s, path):
    # subtle frame behind
    card(s, IMG_L - Inches(0.06), IMG_T - Inches(0.06), IMG_W + Inches(0.12), IMG_H + Inches(0.12), WASH, line=CLOUD)
    cover_img(s, path, IMG_L, IMG_T, IMG_W, IMG_H)

def video_panel(s, video, poster, l, t, w, h, frame=True):
    if frame:
        card(s, l - Inches(0.06), t - Inches(0.06), w + Inches(0.12), h + Inches(0.12), WASH, line=CLOUD)
    return s.shapes.add_movie(video, l, t, w, h, poster_frame_image=poster, mime_type='video/mp4')

VID = os.path.join(ROOT, "assets", "videos")

# SLIDE 2 — THE CHALLENGE
s = slide(); bg(s, WHITE)
header(s, 1, "The Challenge", RED, "A Better Way to Track Attendance.")
txt(s, Inches(0.7), Inches(2.6), Inches(11.9), Inches(0.55),
    "The way schools record attendance today is slow, manual, and error-prone. Here is exactly what the DOS platform replaces — automatically.",
    size=14, color=SLATE, font=FONT, spacing=1.1)
# two comparison cards
cw = Inches(5.9); ch = Inches(3.0); ct = Inches(3.45)
c1 = card(s, Inches(0.7), ct, cw, ch, RGBColor(0xFD,0xF2,0xF2), line=RGBColor(0xF3,0xD0,0xD0))
txt(s, Inches(0.95), Inches(3.6), cw, Inches(0.4), "TODAY", size=12, color=RED, bold=True, font=MONO)
bullets(s, Inches(0.95), Inches(4.05), cw - Inches(0.5), Inches(2.3),
        ["Paper sign-in sheets, prone to errors", "Hours calculated by hand",
         "Payroll delays and inaccuracies", "No real-time oversight for directors",
         "Heavy administrative burden"], size=13, bcolor=RED)
c2 = card(s, Inches(6.75), ct, cw, ch, WASH, line=RGBColor(0xC9,0xDB,0xF5))
txt(s, Inches(7.0), Inches(3.6), cw, Inches(0.4), "WITH DOS", size=12, color=BLUE, bold=True, font=MONO)
bullets(s, Inches(7.0), Inches(4.05), cw - Inches(0.5), Inches(2.3),
        ["Tap-to-check-in from any phone", "Automatic, accurate hour tracking",
         "Payroll-ready exports, instantly", "Real-time director dashboards",
         "RFID-verified & completely paperless"], size=13, bcolor=GREEN)
footer(s)

# SLIDE 3 — THE ECOSYSTEM
s = slide(); bg(s, WHITE)
header(s, 2, "The Ecosystem", BLUE, "One Connected Educational Network.", title_w=Inches(6.7))
txt(s, Inches(0.7), Inches(2.65), TXT_W, Inches(1.2),
    "Every attendance event flows through a secure platform that gives schools, directors, and DOS leadership a shared source of truth across Curaçao's public education system.",
    size=15, color=SLATE, font=FONT, spacing=1.12)
bullets(s, Inches(0.7), Inches(4.1), TXT_W, Inches(2.6),
        [("One teacher", "checks in"), ("One school", "confirms"),
         ("One director", "sees it live"), ("One platform", "updates"),
         ("One payroll", "gets accurate data")], size=15)
img_panel(s, os.path.join(IMG, "pp3-rfid.jpg"))
footer(s)

# SLIDE 4 — HOW IT WORKS
s = slide(); bg(s, WHITE)
header(s, 3, "How It Works", BLUE, "Simple For Teachers. Powerful For Schools.", title_w=Inches(6.7))
txt(s, Inches(0.7), Inches(2.55), TXT_W, Inches(0.9),
    "A secure attendance process that takes only seconds — while automatically updating the entire DOS workforce platform.",
    size=14, color=SLATE, font=FONT, spacing=1.1)
bullets(s, Inches(0.7), Inches(3.5), TXT_W, Inches(3.2),
        ["Teacher arrives at school", "Opens the DOS mobile app",
         "Taps their phone on the RFID reader", "Identity is verified securely",
         "Attendance is recorded instantly", "Director dashboards update automatically",
         "Payroll-ready records are generated"], size=14, gap=6)
# live video (right) + stat chip
video_panel(s, os.path.join(VID, "howitworks.mp4"), os.path.join(IMG, "pp4-checkin.jpg"),
            Inches(7.55), Inches(2.95), Inches(5.08), Inches(2.86))
card(s, Inches(7.55), Inches(6.12), Inches(5.08), Inches(0.62), BLUE, line=BLUE)
txt(s, Inches(7.55), Inches(6.12), Inches(5.08), Inches(0.62),
    "8 seconds  ·  average check-in time", size=14, color=WHITE, bold=True,
    font=FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# SLIDE 5 — LATE ARRIVAL DETECTION
s = slide(); bg(s, WHITE)
header(s, 4, "Automatic Detection", RED, "Late Arrival Detection.", title_w=Inches(6.7))
txt(s, Inches(0.7), Inches(2.65), TXT_W, Inches(1.2),
    "When a teacher checks in late, the platform flags the arrival, calculates the delay, and notifies administration — instantly, with zero manual work.",
    size=15, color=SLATE, font=FONT, spacing=1.12)
bullets(s, Inches(0.7), Inches(4.1), TXT_W, Inches(2.6),
        [("RFID Verified", "— identity confirmed"),
         ("Late Arrival Detected", "— delay calculated live"),
         ("Director Notified", "— push sent instantly"),
         ("Attendance Synced", "— dashboard, real-time"),
         ("Payroll Updated", "— auto-adjusted")], size=15, bcolor=RED)
img_panel(s, os.path.join(IMG, "pp5-late.jpg"))
footer(s)

# SLIDE 6 — MULTI CAMPUS
s = slide(); bg(s, WHITE)
header(s, 5, "Multiple Campus Support", BLUE, "One Teacher. Many Schools.", title_w=Inches(6.7))
txt(s, Inches(0.7), Inches(2.65), TXT_W, Inches(1.2),
    "The platform automatically detects each campus, keeps separate attendance logs, and unifies everything into one DOS report — no manual switching.",
    size=15, color=SLATE, font=FONT, spacing=1.12)
bullets(s, Inches(0.7), Inches(4.1), TXT_W, Inches(2.6),
        [("Automatic Campus Detection", "— knows which school, instantly"),
         ("Separate Attendance Logs", "— each campus kept distinct"),
         ("Unified DOS Reporting", "— one report across every school"),
         ("Real-Time Synchronization", "— updates the moment it happens")], size=15)
img_panel(s, os.path.join(IMG, "pp6-campus.jpg"))
footer(s)

# SLIDE 7 — COMMAND CENTER
s = slide(); bg(s, WHITE)
header(s, 6, "Command Center", BLUE, "Your Entire Network. One View.")
txt(s, Inches(0.7), Inches(2.6), Inches(11.9), Inches(0.55),
    "Every teacher, every school, every attendance event, and every payroll-ready hour across the DOS network — visible from a single command center.",
    size=14, color=SLATE, font=FONT, spacing=1.1)
# stat row
stats = [("12", "Connected Schools"), ("1,349", "Active Staff"),
         ("Real-Time", "Attendance Events"), ("99.9%", "System Uptime")]
sx = Inches(0.7); sw = Inches(2.85); gap = Inches(0.13)
for i,(v,l) in enumerate(stats):
    x = Emu(int(sx) + i*(int(sw)+int(gap)))
    card(s, x, Inches(3.35), sw, Inches(1.25), WASH, line=CLOUD)
    txt(s, x, Inches(3.5), sw, Inches(0.6), v, size=28, color=BLUE, bold=True, font=FONT, align=PP_ALIGN.CENTER)
    txt(s, x, Inches(4.12), sw, Inches(0.4), l, size=11, color=SLATE, font=MONO, align=PP_ALIGN.CENTER)
# dashboard image full width
img = os.path.join(IMG, "Dashboard.png")
card(s, Inches(0.64), Inches(4.84), Inches(12.05), Inches(2.02), WASH, line=CLOUD)
cover_img(s, img, Inches(0.7), Inches(4.9), Inches(11.93), Inches(1.9))
footer(s)

# SLIDE 8 — PAYROLL
s = slide(); bg(s, WHITE)
header(s, 7, "Payroll Preparation", BLUE, "Attendance In. Payroll Out.", title_w=Inches(6.7))
txt(s, Inches(0.7), Inches(2.65), TXT_W, Inches(1.2),
    "Every check-in the platform records automatically becomes a verified, payroll-ready record — no spreadsheets, no manual calculation, no chasing hours.",
    size=15, color=SLATE, font=FONT, spacing=1.12)
bullets(s, Inches(0.7), Inches(4.0), TXT_W, Inches(2.0),
        ["Teacher checks in — time captured", "Synced to the DOS cloud platform",
         "Hours calculated instantly", "One-click SOHO payroll export",
         "Payroll-ready — zero manual work"], size=14)
# comparison chips
card(s, Inches(0.7), Inches(6.15), Inches(3.1), Inches(0.7), RGBColor(0xFD,0xF2,0xF2), line=RGBColor(0xF3,0xD0,0xD0))
txt(s, Inches(0.7), Inches(6.15), Inches(3.1), Inches(0.7), "Manual:  8–12 hours", size=13, color=RED, bold=True, font=FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, Inches(3.95), Inches(6.15), Inches(3.1), Inches(0.7), RGBColor(0xEC,0xFB,0xF4), line=RGBColor(0xBC,0xEAD,0xD8) if False else RGBColor(0xBC,0xEA,0xD8))
txt(s, Inches(3.95), Inches(6.15), Inches(3.1), Inches(0.7), "With DOS:  5 minutes", size=13, color=GREEN, bold=True, font=FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
img_panel(s, os.path.join(IMG, "pp8-payroll.jpg"))
footer(s)

# SLIDE 9 — STRATEGIC VISION (expansion)
s = slide(); bg(s, WHITE)
header(s, 8, "Strategic Vision", BLUE, "School Ecosystem Expansion.")
txt(s, Inches(0.7), Inches(2.55), Inches(11.9), Inches(0.6),
    "Attendance is where it starts — not where it ends. Built on attendance data, the platform grows into the system that runs workforce operations for every DOS school.",
    size=14, color=SLATE, font=FONT, spacing=1.1)
# Today card
card(s, Inches(0.7), Inches(3.35), Inches(3.0), Inches(3.4), BLUE, line=BLUE)
txt(s, Inches(0.9), Inches(3.55), Inches(2.6), Inches(0.4), "TODAY", size=12, color=YELLOW, bold=True, font=MONO)
txt(s, Inches(0.9), Inches(3.95), Inches(2.6), Inches(0.8), "DOS Attendance Platform", size=18, color=WHITE, bold=True, font=FONT, spacing=1.0)
bullets(s, Inches(0.9), Inches(4.95), Inches(2.6), Inches(1.7),
        ["Teacher Check-In", "RFID Verification", "Attendance Tracking", "Director Dashboard"],
        size=12, color=WHITE, bcolor=YELLOW, gap=6)
# Tomorrow grid
mods = ["Payroll Integration", "HR Management", "School Analytics",
        "Multi-School Management", "Leadership Reporting", "Operational Oversight"]
gx = Inches(3.95); gw = Inches(2.78); gh = Inches(1.05); ggap = Inches(0.12)
txt(s, gx, Inches(3.0), Inches(8), Inches(0.35), "TOMORROW", size=12, color=BLUE, bold=True, font=MONO)
for i,m in enumerate(mods):
    r = i // 3; c = i % 3
    x = Emu(int(gx) + c*(int(gw)+int(ggap)))
    y = Emu(int(Inches(3.45)) + r*(int(gh)+int(ggap)) + r*int(Inches(0.45)))
    card(s, x, y, gw, gh, WASH, line=CLOUD)
    txt(s, Emu(int(x)+int(Inches(0.18))), y, Emu(int(gw)-int(Inches(0.3))), gh, m, size=13, color=INK, bold=True, font=FONT, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# SLIDE 10 — SYSTEM ARCHITECTURE (flow diagram)
s = slide(); bg(s, WHITE)
header(s, 9, "System Architecture · Live", BLUE, "Real-Time Data Flow.")
txt(s, Inches(0.7), Inches(2.55), Inches(11.9), Inches(0.55),
    "Watch one attendance tap travel securely through every layer of the platform — from a teacher's phone to payroll — in real time.",
    size=14, color=SLATE, font=FONT, spacing=1.1)
# meta chips
metas = ["AES-256 Encryption", "99.9% Uptime", "0.3s Avg Latency", "Real-Time Sync"]
mx = Inches(0.7)
for i,m in enumerate(metas):
    x = Emu(int(mx) + i*int(Inches(3.0)))
    txt(s, x, Inches(3.25), Inches(2.9), Inches(0.35), "■  " + m, size=11.5, color=BLUE, bold=True, font=MONO)
# flow boxes
steps = [("Teacher Phone","Tap to check in", BLUE, WHITE),
         ("RFID Verification","Secure NFC confirm", YELLOW, INK),
         ("Secure Cloud","Encrypted & stored", BLUE, WHITE),
         ("Director Dashboard","Updates instantly", BLUE, WHITE),
         ("DOS Headquarters","Network oversight", RED, WHITE),
         ("Payroll Export","Payroll-ready", GREEN, WHITE)]
bx = Inches(0.7); bw = Inches(1.79); bh = Inches(1.5); bgap = Inches(0.207); by = Inches(4.2)
for i,(t1,t2,fill,tc) in enumerate(steps):
    x = Emu(int(bx) + i*(int(bw)+int(bgap)))
    card(s, x, by, bw, bh, fill, line=fill)
    txt(s, Emu(int(x)+int(Inches(0.12))), Emu(int(by)+int(Inches(0.22))), Emu(int(bw)-int(Inches(0.24))), Inches(0.75),
        t1, size=13, color=tc, bold=True, font=FONT, align=PP_ALIGN.CENTER, spacing=0.95)
    txt(s, Emu(int(x)+int(Inches(0.1))), Emu(int(by)+int(Inches(0.95))), Emu(int(bw)-int(Inches(0.2))), Inches(0.45),
        t2, size=9.5, color=(RGBColor(0x3A,0x2E,0x00) if fill==YELLOW else RGBColor(0xDD,0xE8,0xF7)), font=MONO, align=PP_ALIGN.CENTER, spacing=0.95)
    if i < len(steps)-1:
        ax = Emu(int(x)+int(bw)+int(Inches(0.02)))
        txt(s, ax, Emu(int(by)+int(Inches(0.5))), Inches(0.2), Inches(0.5), "→", size=18, color=SLATE, bold=True, font=FONT, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.4),
    "Encrypted end-to-end.  Synced in real time.", size=13, color=INK, bold=True, font=FONT)
footer(s)

# SLIDE 11 — PRODUCT EXPERIENCE
s = slide(); bg(s, WHITE)
header(s, 10, "Product Experience", BLUE, "The DOS Platform In Action.")
txt(s, Inches(0.7), Inches(2.55), Inches(11.9), Inches(0.4),
    "Designed for teachers, administrators and school leadership.", size=14, color=SLATE, font=FONT)
# Teacher card — live app video
card(s, Inches(0.7), Inches(3.05), Inches(5.9), Inches(3.85), WHITE, line=CLOUD)
txt(s, Inches(0.95), Inches(3.25), Inches(5.4), Inches(0.4), "Teacher Experience — mobile check-in", size=14, color=INK, bold=True, font=FONT)
video_panel(s, os.path.join(VID, "product.mp4"), os.path.join(IMG, "pp4-checkin.jpg"),
            Inches(0.95), Inches(3.8), Inches(5.4), Inches(2.9), frame=False)
# Admin card — live dashboard video
card(s, Inches(6.75), Inches(3.05), Inches(5.9), Inches(3.85), WHITE, line=CLOUD)
txt(s, Inches(7.0), Inches(3.25), Inches(5.4), Inches(0.4), "Administrator Experience — live dashboard", size=14, color=INK, bold=True, font=FONT)
video_panel(s, os.path.join(VID, "dashboard.mp4"), os.path.join(IMG, "Dashboard.png"),
            Inches(7.0), Inches(3.8), Inches(5.4), Inches(2.9), frame=False)
footer(s)

# SLIDE 12 — IMPLEMENTATION JOURNEY
s = slide(); bg(s, WHITE)
header(s, 11, "Implementation Journey", BLUE, "Rolling Out Across Curaçao.")
txt(s, Inches(0.7), Inches(2.55), Inches(11.9), Inches(0.55),
    "A realistic path from a handful of pilot schools to every public school on the island — the DOS network coming online, stage by stage.",
    size=14, color=SLATE, font=FONT, spacing=1.1)
phases = [("01","Pilot & Foundations","4 schools","First pilot schools, readers installed, staff trained"),
          ("02","Adoption & Impact","10 schools","Teachers adopt tap check-in, accuracy improves"),
          ("03","Insight & Payroll","18 schools","Real-time visibility, automated payroll integration"),
          ("04","Island-Wide","25+ schools","Scaled to every DOS school, one unified platform")]
pw = Inches(2.9); pgap = Inches(0.13); pt = Inches(3.4); ph = Inches(2.6)
for i,(n,t1,sc,d) in enumerate(phases):
    x = Emu(int(Inches(0.7)) + i*(int(pw)+int(pgap)))
    fill = BLUE if i==3 else WASH
    tc = WHITE if i==3 else INK
    dc = RGBColor(0xDD,0xE8,0xF7) if i==3 else SLATE
    card(s, x, pt, pw, ph, fill, line=(BLUE if i==3 else CLOUD))
    txt(s, Emu(int(x)+int(Inches(0.22))), Emu(int(pt)+int(Inches(0.2))), Inches(2.4), Inches(0.4), n, size=15, color=(YELLOW if i==3 else BLUE), bold=True, font=MONO)
    txt(s, Emu(int(x)+int(Inches(0.22))), Emu(int(pt)+int(Inches(0.65))), Emu(int(pw)-int(Inches(0.4))), Inches(0.7), t1, size=16, color=tc, bold=True, font=FONT, spacing=0.95)
    txt(s, Emu(int(x)+int(Inches(0.22))), Emu(int(pt)+int(Inches(1.35))), Emu(int(pw)-int(Inches(0.4))), Inches(0.35), sc, size=12, color=(YELLOW if i==3 else BLUE), bold=True, font=MONO)
    txt(s, Emu(int(x)+int(Inches(0.22))), Emu(int(pt)+int(Inches(1.75))), Emu(int(pw)-int(Inches(0.44))), Inches(0.8), d, size=11, color=dc, font=FONT, spacing=1.05)
txt(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.5),
    "Before DOS: paper & guesswork, 8–12h payroll, delayed data, heavy admin.    With DOS: 100% verified, payroll in minutes, real-time, automated.",
    size=12, color=INK, bold=True, font=FONT, spacing=1.05)
footer(s)

# SLIDE 13 — THE VISION (dark finale)
s = slide(); bg(s, INK)
vis = os.path.join(IMG, "pp13-vision.jpg")
visvid = os.path.join(VID, "vision.mp4")
if os.path.exists(visvid):
    s.shapes.add_movie(visvid, 0, 0, SW, SH, poster_frame_image=vis, mime_type='video/mp4')
    ov = rect(s, 0, 0, SW, SH, INK); set_alpha(ov, 62)
elif os.path.exists(vis):
    cover_img(s, vis, 0, 0, SW, SH)
    ov = rect(s, 0, 0, SW, SH, INK); set_alpha(ov, 62)
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.4),
    "THE VISION  ·  DOS CURAÇAO", size=13, color=YELLOW, bold=True, font=MONO)
txt(s, Inches(0.85), Inches(1.25), Inches(11.6), Inches(1.6),
    "This Is Bigger Than Attendance.", size=46, color=WHITE, bold=True, font=FONT, spacing=0.98)
txt(s, Inches(0.9), Inches(2.95), Inches(9.5), Inches(0.9),
    "DOS is building the digital infrastructure for the future of education across Curaçao — and it starts with a single check-in.",
    size=15, color=RGBColor(0xD8,0xE4,0xF4), font=FONT, spacing=1.12)
manifesto = ["Every School Connected.", "Every Teacher Accounted For.",
             "Every Decision Informed.", "Powered By DOS."]
for i,m in enumerate(manifesto):
    col = YELLOW if i==3 else WHITE
    txt(s, Inches(0.9), Emu(int(Inches(4.0))+i*int(Inches(0.55))), Inches(8), Inches(0.5),
        m, size=20, color=col, bold=True, font=FONT)
# stat strip
vstats = [("25+","Schools Connected"),("100%","Real-Time Visibility"),("One","Unified Platform")]
for i,(v,l) in enumerate(vstats):
    x = Emu(int(Inches(9.0)) + 0)
    y = Emu(int(Inches(4.0)) + i*int(Inches(0.95)))
    txt(s, Inches(9.6), y, Inches(3.2), Inches(0.5), v, size=26, color=YELLOW, bold=True, font=FONT)
    txt(s, Inches(9.6), Emu(int(y)+int(Inches(0.45))), Inches(3.2), Inches(0.35), l, size=10.5, color=RGBColor(0xCF,0xDD,0xF2), font=MONO)
footer(s, dark=True)

# SLIDE 14 — CLOSING
s = slide(); bg(s, INK)
tri_bar(s, 0, 0, SW, Inches(0.18))
lock = os.path.join(LOGO, "dos-lockup.png")
txt(s, Inches(1), Inches(1.5), Inches(11.33), Inches(0.4),
    "REQUEST A DEMO", size=14, color=YELLOW, bold=True, font=MONO, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.05), Inches(11.33), Inches(1.6),
    "DOS isn't just solving attendance.\nIt's building the future of education in Curaçao.",
    size=34, color=WHITE, bold=True, font=FONT, align=PP_ALIGN.CENTER, spacing=1.0)
txt(s, Inches(1), Inches(4.15), Inches(11.33), Inches(0.5),
    "The digital infrastructure for the future of education in Curaçao.",
    size=15, color=RGBColor(0xD8,0xE4,0xF4), font=FONT, align=PP_ALIGN.CENTER)
if os.path.exists(lock):
    contain_img(s, lock, Inches(3.67), Inches(5.0), Inches(6.0), Inches(1.4))
txt(s, Inches(1), Inches(6.85), Inches(11.33), Inches(0.4),
    "© 2026 Dienst Openbare Scholen · Curaçao  ·  Workforce Management Platform · v1.0",
    size=10, color=RGBColor(0x9F,0xB3,0xCC), font=MONO, align=PP_ALIGN.CENTER)

out = os.path.join(ROOT, "DOS-Curacao.pptx")
prs.save(out)
print("SAVED:", out)
print("slides:", len(prs.slides._sldIdLst))
